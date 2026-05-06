from __future__ import annotations

from pathlib import Path

from app.common.schemas import FullCalculationResult
from app.common.utils import utcnow
from app.reporting.charts import save_efficiency_chart, save_temperature_chart


def export_full_result_to_pptx(result: FullCalculationResult, directory: str | Path = "exports") -> Path:
    """Создает краткую PPTX-выгрузку результатов расчета."""
    from pptx import Presentation
    from pptx.util import Inches

    directory = Path(directory)
    directory.mkdir(parents=True, exist_ok=True)
    path = directory / f"calculation_{utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
    chart_dir = directory / "charts"
    heat_chart = save_temperature_chart(result.temperature_analysis, chart_dir / "heat_removal.png")
    eff_chart = save_efficiency_chart(result, chart_dir / "efficiency.png")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Расчет эффективности энергооборудования"
    slide.placeholders[1].text = "Автоматически сформированный отчет программного комплекса"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Основные выходные показатели"
    rows = 8
    cols = 2
    table = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(1.3), Inches(8.8), Inches(3.8)).table
    data = [
        ("Режим работы", result.main_result.operation_mode),
        ("Нагрузка на энергоблок, МВт", result.main_result.load_per_block),
        ("КПД блока, %", result.main_result.block_efficiency * 100),
        ("КПД ТЭС брутто, %", result.main_result.efficiency_brutto * 100),
        ("Собственные нужды, МВт", result.main_result.own_needs_power),
        ("КПД ТЭС нетто, %", result.main_result.efficiency_netto * 100),
        ("Удельный расход топлива, г у.т./кВт·ч", result.main_result.fuel_consumption),
        ("Оптимальное разрежение, кПа", result.condenser_vacuum_optimization.best_vacuum_kpa),
    ]
    for i, (name, value) in enumerate(data):
        table.cell(i, 0).text = name
        if isinstance(value, str):
            table.cell(i, 1).text = value
            continue
        table.cell(i, 1).text = f"{value:.3f}"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Теплоотводящая способность"
    slide.shapes.add_picture(str(heat_chart), Inches(1.0), Inches(1.2), width=Inches(8.0))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "КПД нетто от температуры"
    slide.shapes.add_picture(str(eff_chart), Inches(1.0), Inches(1.2), width=Inches(8.0))

    prs.save(path)
    return path


__all__ = ['export_full_result_to_pptx']
